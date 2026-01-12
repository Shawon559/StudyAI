"""
RAG Ingestion - Process files and YouTube for RAG pipeline
"""
import os
import uuid
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation
from docx import Document
from pypdf import PdfReader
import whisper

RAW_DIR = "rag/data/raw"

print("Loading Whisper model (base)...")
whisper_model = whisper.load_model("base")


def save_text(text, lecture_id=None):
    if not lecture_id:
        lecture_id = str(uuid.uuid4())[:8]
    out_path = os.path.join(RAW_DIR, f"{lecture_id}.txt")
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(text)
    print(f"[OK] Saved → {out_path}")
    return out_path


def ingest_pdf(path):
    reader = PdfReader(path)
    return "\n".join([p.extract_text() or "" for p in reader.pages])


def ingest_ppt(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def ingest_docx(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])


def ingest_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def ingest_youtube(url):
    video_id = None
    if "v=" in url:
        video_id = url.split("v=")[-1].split("&")[0]
    elif "youtu.be/" in url:
        video_id = url.split("youtu.be/")[-1].split("?")[0]

    if not video_id:
        raise ValueError("Invalid YouTube URL")

    api = YouTubeTranscriptApi()
    fetched = api.fetch(video_id)
    return " ".join(snippet.text for snippet in fetched.snippets)


def transcribe_audio(path):
    print(f"[Whisper] Transcribing {path}...")
    result = whisper_model.transcribe(path)
    return result["text"]


def ingest_video(path):
    from moviepy.video.io.VideoFileClip import VideoFileClip
    print("[Video] Extracting audio...")
    clip = VideoFileClip(path)
    audio_path = path + "_temp_audio.wav"
    clip.audio.write_audiofile(audio_path)
    text = transcribe_audio(audio_path)
    os.remove(audio_path)
    return text


def ingest(path_or_url, lecture_id=None):
    """Ingest any file or YouTube link"""
    if lecture_id is None:
        lecture_id = str(uuid.uuid4())[:8]

    if path_or_url.startswith("http"):
        text = ingest_youtube(path_or_url)
        return save_text(text, lecture_id)

    path = path_or_url.lower()

    if path.endswith(".pdf"):
        text = ingest_pdf(path_or_url)
    elif path.endswith((".ppt", ".pptx")):
        text = ingest_ppt(path_or_url)
    elif path.endswith((".doc", ".docx")):
        text = ingest_docx(path_or_url)
    elif path.endswith(".txt"):
        text = ingest_txt(path_or_url)
    elif path.endswith((".mp3", ".wav", ".m4a")):
        text = transcribe_audio(path_or_url)
    elif path.endswith((".mp4", ".mkv", ".mov", ".avi")):
        text = ingest_video(path_or_url)
    else:
        raise ValueError("Unsupported file type: " + path)

    return save_text(text, lecture_id)


if __name__ == "__main__":
    path = input("Enter file path or YouTube URL: ")
    ingest(path)

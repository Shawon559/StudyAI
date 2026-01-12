"""
Multi-Modal Processor - Handles text, audio, video, images, documents
"""
import os
import tempfile
from typing import Dict, Union


class MultiModalProcessor:
    """Process various input types into text"""

    def __init__(self):
        self.supported_types = {
            'text': ['.txt', '.md', '.json', '.csv'],
            'audio': ['.mp3', '.wav', '.m4a', '.ogg', '.flac'],
            'image': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'],
            'video': ['.mp4', '.avi', '.mov', '.mkv', '.webm'],
            'document': ['.pdf', '.doc', '.docx', '.ppt', '.pptx']
        }

    def detect_input_type(self, filename: str = None, mime_type: str = None) -> str:
        """Detect input modality"""
        if filename:
            ext = os.path.splitext(filename)[1].lower()
            for modality, extensions in self.supported_types.items():
                if ext in extensions:
                    return modality

        if mime_type:
            if mime_type.startswith('text/'): return 'text'
            elif mime_type.startswith('audio/'): return 'audio'
            elif mime_type.startswith('image/'): return 'image'
            elif mime_type.startswith('video/'): return 'video'
            elif 'pdf' in mime_type or 'document' in mime_type: return 'document'

        return 'unknown'

    def process_text(self, text: str) -> Dict:
        """Process plain text"""
        return {'modality': 'text', 'content': text, 'length': len(text), 'status': 'success'}

    def process_audio(self, audio_file_path: str) -> Dict:
        """Transcribe audio using Whisper (local, free)"""
        try:
            import whisper
            import torch

            if torch.backends.mps.is_available(): device = "mps"
            elif torch.cuda.is_available(): device = "cuda"
            else: device = "cpu"

            print(f"[Whisper] Loading on {device}...")
            model = whisper.load_model("base", device=device)

            print(f"[Whisper] Transcribing {audio_file_path}...")
            result = model.transcribe(audio_file_path)

            return {
                'modality': 'audio', 'content': result["text"],
                'length': len(result["text"]), 'status': 'success',
                'original_file': audio_file_path
            }
        except Exception as e:
            return {'modality': 'audio', 'content': '', 'status': 'error', 'error': str(e)}

    def process_image(self, image_file_path: str, prompt: str = None) -> Dict:
        """Process image using GPT-4 Vision"""
        try:
            import base64
            from openai import OpenAI

            client = OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

            with open(image_file_path, 'rb') as f:
                image_data = base64.b64encode(f.read()).decode('utf-8')

            ext = os.path.splitext(image_file_path)[1].lower()
            mime_type = "image/jpeg" if ext == '.jpg' else f"image/{ext[1:]}"

            vision_prompt = prompt or "Extract all text and describe educational content in this image."

            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": vision_prompt},
                        {"type": "image_url", "image_url": {"url": f"data:{mime_type};base64,{image_data}"}}
                    ]
                }],
                max_tokens=1000
            )

            content = response.choices[0].message.content
            return {
                'modality': 'image', 'content': content,
                'length': len(content), 'status': 'success',
                'original_file': image_file_path
            }
        except Exception as e:
            return {'modality': 'image', 'content': '', 'status': 'error', 'error': str(e)}

    def process_video(self, video_file_path: str) -> Dict:
        """Process video - extract audio and transcribe"""
        try:
            import subprocess
            import shutil

            temp_dir = tempfile.mkdtemp()
            audio_path = os.path.join(temp_dir, 'audio.mp3')

            # Extract audio
            try:
                subprocess.run([
                    'ffmpeg', '-i', video_file_path,
                    '-q:a', '0', '-map', 'a', audio_path, '-y'
                ], check=True, capture_output=True)

                audio_result = self.process_audio(audio_path)
                content = audio_result.get('content', '')
            except Exception as e:
                content = f"[Audio extraction failed: {e}]"

            shutil.rmtree(temp_dir, ignore_errors=True)

            return {
                'modality': 'video', 'content': content,
                'length': len(content), 'status': 'success',
                'original_file': video_file_path
            }
        except Exception as e:
            return {'modality': 'video', 'content': '', 'status': 'error', 'error': str(e)}

    def process_document(self, document_path: str) -> Dict:
        """Process PDF, Word, PowerPoint"""
        try:
            ext = os.path.splitext(document_path)[1].lower()

            # Word documents
            if ext in ['.doc', '.docx']:
                from docx import Document
                doc = Document(document_path)
                text = '\n'.join([p.text for p in doc.paragraphs])
                return {'modality': 'document', 'content': text, 'length': len(text), 'status': 'success'}

            # PowerPoint
            if ext in ['.ppt', '.pptx']:
                from pptx import Presentation
                prs = Presentation(document_path)
                text = '\n'.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                return {'modality': 'document', 'content': text, 'length': len(text), 'status': 'success'}

            # PDF - try RAG ingest first
            try:
                from rag.ingest import ingest_file
                result = ingest_file(document_path)
                return {'modality': 'document', 'content': result.get('text', ''), 'status': 'success'}
            except:
                pass

            # Fallback PDF
            from pypdf import PdfReader
            reader = PdfReader(document_path)
            text = '\n'.join([p.extract_text() or '' for p in reader.pages])
            return {'modality': 'document', 'content': text, 'length': len(text), 'status': 'success'}

        except Exception as e:
            return {'modality': 'document', 'content': '', 'status': 'error', 'error': str(e)}

    def process_youtube(self, youtube_url: str) -> Dict:
        """Process YouTube - get transcript or transcribe audio"""
        import re

        # Extract video ID
        patterns = [
            r'(?:youtube\.com\/watch\?v=|youtu\.be\/|youtube\.com\/embed\/)([^&\n?#]+)',
            r'youtube\.com\/shorts\/([^&\n?#]+)'
        ]

        video_id = None
        for pattern in patterns:
            match = re.search(pattern, youtube_url)
            if match:
                video_id = match.group(1)
                break

        if not video_id:
            return {'modality': 'video', 'content': '', 'status': 'error', 'error': 'Invalid YouTube URL'}

        # Try subtitles first (free, fast)
        try:
            from youtube_transcript_api import YouTubeTranscriptApi
            transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
            text = ' '.join([item['text'] for item in transcript_list])
            return {
                'modality': 'video', 'content': text, 'length': len(text),
                'status': 'success', 'video_id': video_id, 'source': 'youtube_subtitles'
            }
        except Exception as e:
            print(f"[YouTube] No subtitles: {e}, trying audio...")

        # Fallback: download and transcribe
        try:
            import subprocess
            import glob
            import shutil

            temp_dir = tempfile.mkdtemp()
            audio_path = os.path.join(temp_dir, 'youtube_audio.mp3')

            # Download with yt-dlp
            subprocess.run([
                'yt-dlp', '-x', '--audio-format', 'mp3',
                '-o', audio_path.replace('.mp3', '.%(ext)s'), youtube_url
            ], capture_output=True, timeout=120)

            audio_files = glob.glob(os.path.join(temp_dir, 'youtube_audio.*'))
            if audio_files:
                actual_audio = audio_files[0]
                if not actual_audio.endswith('.mp3'):
                    subprocess.run(['ffmpeg', '-i', actual_audio, '-q:a', '0', audio_path, '-y'], capture_output=True)
                else:
                    audio_path = actual_audio

                audio_result = self.process_audio(audio_path)
                shutil.rmtree(temp_dir, ignore_errors=True)

                if audio_result['status'] == 'success':
                    return {
                        'modality': 'video', 'content': audio_result['content'],
                        'length': len(audio_result['content']), 'status': 'success',
                        'video_id': video_id, 'source': 'whisper_transcription'
                    }

            shutil.rmtree(temp_dir, ignore_errors=True)
            return {'modality': 'video', 'content': '', 'status': 'error', 'error': 'Audio download failed'}

        except Exception as e:
            return {'modality': 'video', 'content': '', 'status': 'error', 'error': str(e)}

    def process(self, input_data: Union[str, Dict], input_type: str = None) -> Dict:
        """Main processor - routes to appropriate handler"""
        # Text input
        if isinstance(input_data, str) and not os.path.exists(input_data):
            return self.process_text(input_data)

        # File path
        if isinstance(input_data, str) and os.path.exists(input_data):
            file_path = input_data
            if not input_type:
                input_type = self.detect_input_type(filename=file_path)
        else:
            return {'modality': 'unknown', 'content': '', 'status': 'error', 'message': 'Invalid input'}

        # Route to processor
        if input_type == 'text':
            with open(file_path, 'r', encoding='utf-8') as f:
                return self.process_text(f.read())
        elif input_type == 'audio':
            return self.process_audio(file_path)
        elif input_type == 'image':
            return self.process_image(file_path)
        elif input_type == 'video':
            return self.process_video(file_path)
        elif input_type == 'document':
            return self.process_document(file_path)

        return {'modality': 'unknown', 'content': '', 'status': 'error', 'message': f'Unsupported type: {input_type}'}
